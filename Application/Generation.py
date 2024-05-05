import ast
from pptx import Presentation
from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
import io

def generate_powerpoint(Results, company_name):
    for result in Results:
        if str(list(result.keys())[0]) != 'Competitor' and str(list(result.keys())[0]) != 'Diligence':
            start_index = str(list(result.values())[0]).find('{')
            end_index = str(list(result.values())[0]).rfind('}')
            json_string = str(list(result.values())[0])[start_index:end_index + 1]
            try:
                data = ast.literal_eval(json_string)
                result[str(list(result.keys())[0])] = data
            except (SyntaxError, ValueError) as e:
                print("无法将字符串转换为字典:", e)
        elif str(list(result.keys())[0]) == 'Competitor' or str(list(result.keys())[0]) == 'Diligence':
            llm = ChatOpenAI(model="gpt-4-turbo-preview", temperature=0.0)
            prompt = PromptTemplate(
                template="""Read the following text, split it into parts based on the key points it has, return the parts in a python list. Keep the original texts.
            You should not include any unneccessary information such as introduction or conclusion as key points.
            {result}
            """,
                input_variables=["result"]
            )
            prompt_and_model = prompt | llm
            output = prompt_and_model.invoke({"result": str(list(result.values())[0])})
            start_index = output.content.find('[')
            end_index = output.content.rfind(']')
            list_string = output.content[start_index:end_index + 1]
            print(ast.literal_eval(list_string))
            result[str(list(result.keys())[0])] = ast.literal_eval(list_string)

    prs = Presentation('template.pptx')
    page_to_be_deleted = []
    slide = prs.slides[0]
    text_to_replace = {
        "{company}": company_name
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for key, value in text_to_replace.items():
                    if key in paragraph.text:
                        paragraph.text = paragraph.text.replace(key, value)
    for item in Results:
        if str(list(item.keys())[0]) == 'Introduction':
            result = item['Introduction']
            slide = prs.slides[1]
            text_to_replace = {
                "{c}": result['Company Name'],
                "{s}": result['Sector'],
                "{i}": result['Industry'],
                "{co}": result['Country'],
                "{ci}": result['City'],
                "{ee}": result['Number of employees'],
                "{w}": result['Website'],
                "{summary}": result['Overall Introduction']
            }
        elif str(list(item.keys())[0]) == 'SWOT':
            result = item['SWOT']
            slide = prs.slides[2]
            text_to_replace = {
                "{SWOTtitle}": "SWOT Analysis",
                "{s}": result['Strengths'],
                "{w}": result['Weaknesses'],
                "{o}": result['Opportunities'],
                "{t}": result['Threats']
            }
        elif str(list(item.keys())[0]) == 'Porter':
            result = item['Porter']
            slide = prs.slides[3]
            text_to_replace = {
                "{porters_title}": "Porter Five Forces Analysis",
                "{bargainingS}": result['Bargaining Power of Suppliers'],
                "{bargainingB}": result['Bargaining Power of Buyers'],
                "{intensity}": result['Intensity of Competitive Rivalry'],
                "{threatsN}": result['Threats of New Entrants'],
                "{threatsS}": result['Threats of Substitute Products']
            }
        elif str(list(item.keys())[0]) == 'Canvas':
            result = item['Canvas']
            slide = prs.slides[4]
            text_to_replace = {
                "{bmc_title}": "Business Canvas",
                "{key_partnerships}": result['Key Partners'],
                "{key_activities}": result['Key Activities'],
                "{key_resources}": result['Key Resources'],
                "{value_propositions}": result['Value Proposition'],
                "{customer_relationships}": result['Customer Relationships'],
                "{channels}": result['Channels'],
                "{customer_segments}": result['Customer Segments'],
                "{cost_structure}": result['Cost Structure'],
                "{revenue_streams}": result['Revenue Stream']
            }
        elif str(list(item.keys())[0]) == 'Ansoff':
            result = item['Ansoff']
            slide = prs.slides[5]
            text_to_replace = {
                "{ansoff_title}": "Ansoff Matrix Analysis",
                "{market_development_strategy}": result['Market Development'],
                "{diversification_strategy}": result['Diversification'],
                "{market_penetration_strategy}": result['Market Penetration'],
                "{product_development_strategy}": result['Product Development']
            }
        elif str(list(item.keys())[0]) == 'Competitor':
            result = item['Competitor']
            page_num = len(result) + 4
            slide = prs.slides[page_num]
            if len(result) == 2:
                text_to_replace = {
                    "{Title}": "Competitor Analysis",
                    "{1}": result[0],
                    "{2}": result[1]
                }
                page_to_be_deleted.extend([7, 8, 9])
            elif len(result) == 3:
                text_to_replace = {
                    "{Title}": "Competitor Analysis",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2]
                }
                page_to_be_deleted.extend([6, 8, 9])
            elif len(result) == 4:
                text_to_replace = {
                    "{Title}": "Competitor Analysis",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2],
                    "{4}": result[3]
                }
                page_to_be_deleted.extend([6, 7, 9])
            elif len(result) == 5:
                text_to_replace = {
                    "{Title}": "Competitor Analysis",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2],
                    "{4}": result[3],
                    "{5}": result[4]
                }
                page_to_be_deleted.extend([6, 7, 8])
        elif str(list(item.keys())[0]) == 'Diligence':
            result = item['Diligence']
            page_num = len(result) + 8
            slide = prs.slides[page_num]
            if len(result) == 2:
                text_to_replace = {
                    "{Title}": "Due Diligence",
                    "{1}": result[0],
                    "{2}": result[1]
                }
                page_to_be_deleted.extend([11, 12, 13])
            elif len(result) == 3:
                text_to_replace = {
                    "{Title}": "Due Diligence",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2]
                }
                page_to_be_deleted.extend([10, 12, 13])
            elif len(result) == 4:
                text_to_replace = {
                    "{Title}": "Due Diligence",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2],
                    "{4}": result[3]
                }
                page_to_be_deleted.extend([10, 11, 13])
            elif len(result) == 5:
                text_to_replace = {
                    "{Title}": "Due Diligence",
                    "{1}": result[0],
                    "{2}": result[1],
                    "{3}": result[2],
                    "{4}": result[3],
                    "{5}": result[4]
                }
                page_to_be_deleted.extend([10, 11, 12])

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for key, value in text_to_replace.items():
                        if key in paragraph.text:
                            paragraph.text = paragraph.text.replace(key, value)

    sldids = prs.slides._sldIdLst  # 创建slideid列表，每个id对应一张ppt
    slides = prs.slides  # ppt对象的slides属性，对应每页ppt构成的集合
    index_to_be_deleted = sorted(page_to_be_deleted, reverse=True)
    for index in index_to_be_deleted:  # 删除第一张和第二张，这个对应的是python列表的操作方法
        slides._sldIdLst.remove(sldids[index])
    prs.save('output.pptx')
